# -*- coding: utf-8 -*-
"""Private DOCX/PPTX delivery for completed JEAC reports.

The public repository never receives the generated documents.  A GitHub Actions
runner creates them in a temporary directory, then uploads them only to the
Google Drive folder explicitly shared with the Service Account.
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime, timezone, timedelta
from pathlib import Path
from typing import Iterable, List, Optional
import json
import logging
import os
import re
import tempfile

logger = logging.getLogger(__name__)

# The service account itself is shared only with the designated report folder.
# Full Drive scope is required to create and find subfolders under that shared
# parent; it does not grant access to the owner's unshared Drive content.
_DRIVE_SCOPE = "https://www.googleapis.com/auth/drive"
_TAIPEI = timezone(timedelta(hours=8))
_SUPPORTED_KINDS = {"daily", "weekly", "monthly"}


class PrivateReportDeliveryError(RuntimeError):
    """Raised when enabled private delivery cannot safely complete."""


@dataclass(frozen=True)
class DeliveredReport:
    report_kind: str
    format: str
    file_name: str
    drive_file_id: str
    drive_url: Optional[str] = None


def _env_bool(name: str, default: bool = False) -> bool:
    value = os.getenv(name, "")
    if not value:
        return default
    return value.strip().lower() in {"1", "true", "yes", "on"}


def _markdown_text(value: str) -> str:
    text = re.sub(r"^#{1,6}\s*", "", value.strip())
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"\[(.*?)\]\([^)]*\)", r"\1", text)
    return text.strip()


def _safe_segment(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9_-]+", "_", value.strip()) or "Report"


class PrivateReportDelivery:
    """Generate and upload completed reports without exposing them as artifacts."""

    def __init__(
        self,
        *,
        enabled: bool,
        folder_id: str = "",
        service_account_json: str = "",
        output_dir: Optional[Path] = None,
    ) -> None:
        self.enabled = enabled
        self.folder_id = folder_id.strip()
        self.service_account_json = service_account_json.strip()
        self.output_dir = output_dir

    @classmethod
    def from_environment(cls) -> "PrivateReportDelivery":
        return cls(
            enabled=_env_bool("PRIVATE_REPORT_EXPORT_ENABLED", False),
            folder_id=os.getenv("GOOGLE_DRIVE_REPORT_FOLDER_ID", ""),
            service_account_json=os.getenv("GOOGLE_DRIVE_SERVICE_ACCOUNT_JSON", ""),
        )

    def export(
        self,
        *,
        report_kind: str,
        markdown: str,
        report_date: Optional[datetime] = None,
    ) -> List[DeliveredReport]:
        """Create DOCX and, for weekly/monthly, PPTX only after integrity passes."""
        if not self.enabled:
            logger.info("[private-report] export disabled")
            return []

        kind = str(report_kind or "").strip().lower()
        if kind not in _SUPPORTED_KINDS:
            raise PrivateReportDeliveryError(f"Unsupported private report kind: {report_kind}")
        if not markdown or not markdown.strip():
            raise PrivateReportDeliveryError("Refusing to export an empty report")
        if not self.folder_id or not self.service_account_json:
            raise PrivateReportDeliveryError(
                "Private report export is enabled but Google Drive credentials or folder ID are missing"
            )

        moment = report_date.astimezone(_TAIPEI) if report_date else datetime.now(_TAIPEI)
        date_label = moment.strftime("%Y-%m") if kind == "monthly" else moment.strftime("%Y-%m-%d")
        title = f"JEAC {kind.title()} Investment Report"
        base_name = f"JEAC_{kind.title()}_{date_label}"

        with tempfile.TemporaryDirectory(prefix="jeac-private-report-") as temporary:
            directory = self.output_dir or Path(temporary)
            directory.mkdir(parents=True, exist_ok=True)
            docx_path = directory / f"{base_name}.docx"
            self._create_docx(docx_path, title=title, markdown=markdown, generated_at=moment)

            uploads = [(docx_path, "docx")]
            if kind in {"weekly", "monthly"}:
                pptx_path = directory / f"{base_name}.pptx"
                self._create_pptx(pptx_path, title=title, markdown=markdown, generated_at=moment)
                uploads.append((pptx_path, "pptx"))

            drive = self._build_drive_client()
            kind_folder = self._ensure_folder(drive, self.folder_id, kind.title())
            delivered = [
                self._upload_file(drive, item_path, kind=kind, format_name=format_name, parent_id=kind_folder)
                for item_path, format_name in uploads
            ]

        logger.info(
            "[private-report] uploaded %d private %s document(s) to Google Drive",
            len(delivered),
            kind,
        )
        return delivered

    @staticmethod
    def _create_docx(path: Path, *, title: str, markdown: str, generated_at: datetime) -> None:
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.shared import Inches, Pt, RGBColor
        except ImportError as exc:  # pragma: no cover - exercised in deployment checks
            raise PrivateReportDeliveryError("python-docx is not installed") from exc

        document = Document()
        section = document.sections[0]
        section.top_margin = Inches(0.75)
        section.bottom_margin = Inches(0.75)
        section.left_margin = Inches(0.8)
        section.right_margin = Inches(0.8)

        normal = document.styles["Normal"]
        normal.font.name = "Microsoft JhengHei"
        normal.font.size = Pt(10.5)
        normal.paragraph_format.space_after = Pt(6)

        heading_colours = {
            1: RGBColor(31, 78, 121),
            2: RGBColor(47, 84, 150),
            3: RGBColor(89, 89, 89),
        }
        for level, colour in heading_colours.items():
            style = document.styles[f"Heading {level}"]
            style.font.name = "Microsoft JhengHei"
            style.font.color.rgb = colour

        heading = document.add_heading(title, level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle = document.add_paragraph(
            f"資料完整性通過後產生｜{generated_at.strftime('%Y-%m-%d %H:%M')}（Asia/Taipei）"
        )
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle.runs[0].font.size = Pt(9)
        document.add_paragraph("")

        lines = markdown.splitlines()
        index = 0
        while index < len(lines):
            line = lines[index].rstrip()
            if not line.strip() or line.strip() == "---":
                index += 1
                continue
            if line.lstrip().startswith("|"):
                table_lines: List[str] = []
                while index < len(lines) and lines[index].lstrip().startswith("|"):
                    table_lines.append(lines[index])
                    index += 1
                PrivateReportDelivery._append_docx_table(document, table_lines)
                continue
            heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading_match:
                level = min(len(heading_match.group(1)), 3)
                document.add_heading(_markdown_text(heading_match.group(2)), level=level)
            elif re.match(r"^\s*[-*+]\s+", line):
                document.add_paragraph(_markdown_text(re.sub(r"^\s*[-*+]\s+", "", line)), style="List Bullet")
            elif re.match(r"^\s*\d+[.)]\s+", line):
                document.add_paragraph(
                    _markdown_text(re.sub(r"^\s*\d+[.)]\s+", "", line)),
                    style="List Number",
                )
            else:
                document.add_paragraph(_markdown_text(line))
            index += 1

        document.save(path)

    @staticmethod
    def _append_docx_table(document, lines: Iterable[str]) -> None:
        rows = [
            [_markdown_text(cell) for cell in raw.strip().strip("|").split("|")]
            for raw in lines
            if raw.strip() and not re.fullmatch(r"\|?\s*:?-{3,}:?\s*(\|\s*:?-{3,}:?\s*)+\|?", raw.strip())
        ]
        if not rows:
            return
        column_count = max(len(row) for row in rows)
        table = document.add_table(rows=1, cols=column_count)
        table.style = "Table Grid"
        for row_index, row_values in enumerate(rows):
            cells = table.rows[0].cells if row_index == 0 else table.add_row().cells
            for column_index, value in enumerate(row_values):
                cells[column_index].text = value
        document.add_paragraph("")

    @staticmethod
    def _create_pptx(path: Path, *, title: str, markdown: str, generated_at: datetime) -> None:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as exc:  # pragma: no cover - exercised in deployment checks
            raise PrivateReportDeliveryError("python-pptx is not installed") from exc

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title_box = title_slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.name = "Microsoft JhengHei"
        title_paragraph.font.size = Pt(38)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(31, 78, 121)
        title_paragraph.alignment = PP_ALIGN.CENTER

        date_box = title_slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = f"資料完整性通過後產生｜{generated_at.strftime('%Y-%m-%d %H:%M')}（Asia/Taipei）"
        date_paragraph = date_frame.paragraphs[0]
        date_paragraph.font.name = "Microsoft JhengHei"
        date_paragraph.font.size = Pt(16)
        date_paragraph.alignment = PP_ALIGN.CENTER

        for section_title, section_lines in PrivateReportDelivery._markdown_sections(markdown)[:7]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            banner = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(0.65))
            banner_frame = banner.text_frame
            banner_frame.text = section_title
            banner_paragraph = banner_frame.paragraphs[0]
            banner_paragraph.font.name = "Microsoft JhengHei"
            banner_paragraph.font.size = Pt(28)
            banner_paragraph.font.bold = True
            banner_paragraph.font.color.rgb = RGBColor(31, 78, 121)

            body = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.8), Inches(5.5))
            frame = body.text_frame
            frame.clear()
            frame.word_wrap = True
            for line_number, line in enumerate(section_lines[:10]):
                paragraph = frame.paragraphs[0] if line_number == 0 else frame.add_paragraph()
                paragraph.text = line
                paragraph.level = 0
                paragraph.font.name = "Microsoft JhengHei"
                paragraph.font.size = Pt(19)
                paragraph.space_after = Pt(9)
        presentation.save(path)

    @staticmethod
    def _markdown_sections(markdown: str) -> List[tuple[str, List[str]]]:
        sections: List[tuple[str, List[str]]] = []
        current_title = "投資策略摘要"
        current_lines: List[str] = []
        for raw in markdown.splitlines():
            line = raw.strip()
            heading = re.match(r"^#{1,3}\s+(.+)$", line)
            if heading:
                if current_lines:
                    sections.append((current_title, current_lines))
                current_title = _markdown_text(heading.group(1))
                current_lines = []
                continue
            if not line or line == "---":
                continue
            if line.startswith("|"):
                cells = [_markdown_text(cell) for cell in line.strip("|").split("|")]
                if not all(re.fullmatch(r"\s*:?-{3,}:?\s*", cell) for cell in cells):
                    current_lines.append("｜".join(cells))
                continue
            current_lines.append(_markdown_text(re.sub(r"^[-*+]\s+", "", line)))
        if current_lines:
            sections.append((current_title, current_lines))
        return sections or [("投資策略摘要", ["報告未提供可轉換的內容。"])]

    def _build_drive_client(self):
        try:
            from google.oauth2 import service_account
            from googleapiclient.discovery import build
        except ImportError as exc:  # pragma: no cover - exercised in deployment checks
            raise PrivateReportDeliveryError("Google Drive client dependencies are not installed") from exc
        try:
            payload = json.loads(self.service_account_json)
        except json.JSONDecodeError as exc:
            raise PrivateReportDeliveryError("GOOGLE_DRIVE_SERVICE_ACCOUNT_JSON is not valid JSON") from exc
        credentials = service_account.Credentials.from_service_account_info(payload, scopes=[_DRIVE_SCOPE])
        return build("drive", "v3", credentials=credentials, cache_discovery=False)

    @staticmethod
    def _ensure_folder(drive, parent_id: str, folder_name: str) -> str:
        escaped_name = folder_name.replace("'", "\\'")
        response = drive.files().list(
            q=(
                f"name = '{escaped_name}' and mimeType = 'application/vnd.google-apps.folder' "
                f"and '{parent_id}' in parents and trashed = false"
            ),
            spaces="drive",
            fields="files(id,name)",
            supportsAllDrives=True,
            includeItemsFromAllDrives=True,
        ).execute()
        matches = response.get("files", [])
        if matches:
            return str(matches[0]["id"])
        created = drive.files().create(
            body={
                "name": _safe_segment(folder_name),
                "mimeType": "application/vnd.google-apps.folder",
                "parents": [parent_id],
            },
            fields="id",
            supportsAllDrives=True,
        ).execute()
        return str(created["id"])

    @staticmethod
    def _upload_file(drive, path: Path, *, kind: str, format_name: str, parent_id: str) -> DeliveredReport:
        try:
            from googleapiclient.http import MediaFileUpload
        except ImportError as exc:  # pragma: no cover
            raise PrivateReportDeliveryError("Google Drive upload dependency is not installed") from exc

        mime_type = (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            if format_name == "docx"
            else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        response = drive.files().create(
            body={"name": path.name, "parents": [parent_id]},
            media_body=MediaFileUpload(str(path), mimetype=mime_type, resumable=False),
            fields="id,webViewLink",
            supportsAllDrives=True,
        ).execute()
        return DeliveredReport(
            report_kind=kind,
            format=format_name,
            file_name=path.name,
            drive_file_id=str(response["id"]),
            drive_url=response.get("webViewLink"),
        )
